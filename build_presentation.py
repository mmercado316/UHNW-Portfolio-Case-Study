#!/usr/bin/env python3
"""
UHNW Portfolio Transition - PowerPoint Presentation Builder
Run:    python3 build_presentation.py
Output: UHNW_Portfolio_Presentation.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── PATHS ─────────────────────────────────────────────────────────────────────
BASE_DIR    = os.path.dirname(os.path.abspath(__file__))
OUTPUT_FILE = os.path.join(BASE_DIR, "UHNW_Portfolio_Presentation.pptx")

# ── COLORS ────────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x00, 0x2D, 0x72)
GOLD   = RGBColor(0xC8, 0xA8, 0x4B)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF2, 0xF2, 0xF2)
DGRAY  = RGBColor(0x40, 0x40, 0x40)
MGRAY  = RGBColor(0xA0, 0xA0, 0xA0)
GREEN  = RGBColor(0x00, 0x70, 0x3C)
RED    = RGBColor(0xC0, 0x00, 0x00)
LBLUE  = RGBColor(0xE8, 0xEE, 0xF7)

# ── SLIDE DIMENSIONS (Widescreen 16:9) ────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # truly blank layout

# ══════════════════════════════════════════════════════════════════════════════
# HELPER FUNCTIONS
# ══════════════════════════════════════════════════════════════════════════════

def add_slide():
    return prs.slides.add_slide(BLANK)

def rgb(r, g, b):
    return RGBColor(r, g, b)

def rect(slide, x, y, w, h, fill_color, alpha=None):
    """Add a filled rectangle."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def textbox(slide, x, y, w, h, text, size=18, bold=False, color=WHITE,
            align=PP_ALIGN.LEFT, italic=False, wrap=True, line_spacing=None):
    """Add a text box."""
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name   = "Calibri"
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def multiline_textbox(slide, x, y, w, h, lines, size=14, bold=False,
                      color=WHITE, align=PP_ALIGN.LEFT, spacing_before=0):
    """Add a text box with multiple paragraphs."""
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, (text, fsize, fbold, fcolor) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if spacing_before and i > 0:
            p.space_before = Pt(spacing_before)
        run = p.add_run()
        run.text = text
        run.font.name  = "Calibri"
        run.font.size  = Pt(fsize)
        run.font.bold  = fbold
        run.font.color.rgb = fcolor
    return txb

def nav_bar(slide, current_num):
    """Bottom navigation bar showing slide context."""
    rect(slide, 0, 7.1, 13.33, 0.4, NAVY)
    sections = [
        (1,  "Cover"),
        (2,  "Agenda"),
        (3,  "Challenge"),
        (4,  "Risk Data"),
        (5,  "Recommendation"),
        (6,  "Allocation"),
        (7,  "Tax Strategy"),
        (8,  "SBL Credit"),
        (9,  "Goals"),
        (10, "Timeline"),
        (11, "Summary"),
    ]
    step = 13.33 / len(sections)
    for num, label in sections:
        x = (num - 1) * step
        c = GOLD if num == current_num else MGRAY
        # size bumped 7→9 for readability
        textbox(slide, x + 0.05, 7.13, step - 0.1, 0.34,
                label, size=9, color=c,
                bold=(num == current_num), align=PP_ALIGN.CENTER)

def slide_header(slide, title, subtitle=None, accent_bar=True):
    """Standard slide header with optional gold accent bar."""
    rect(slide, 0, 0, 13.33, 1.1, NAVY)
    if accent_bar:
        rect(slide, 0, 1.1, 13.33, 0.06, GOLD)
    # title bumped 24→26, subtitle 13→14
    textbox(slide, 0.4, 0.08, 12.5, 0.68,
            title, size=26, bold=True, color=WHITE)
    if subtitle:
        textbox(slide, 0.4, 0.70, 12.5, 0.36,
                subtitle, size=14, color=GOLD)

def kpi_box(slide, x, y, w, h, label, value, sub=None,
            val_color=WHITE, bg=NAVY):
    """Metric KPI card."""
    rect(slide, x, y, w, h, bg)
    # label bumped 9→10, sub bumped 8→10
    textbox(slide, x + 0.1, y + 0.08, w - 0.2, 0.3,
            label, size=10, color=MGRAY, align=PP_ALIGN.CENTER)
    textbox(slide, x + 0.05, y + 0.35, w - 0.1, 0.55,
            value, size=22, bold=True, color=val_color,
            align=PP_ALIGN.CENTER)
    if sub:
        textbox(slide, x + 0.05, y + 0.88, w - 0.1, 0.25,
                sub, size=10, color=MGRAY, align=PP_ALIGN.CENTER)

def bullet_block(slide, x, y, w, h, items, size=13, color=DGRAY,
                 bullet="▪", indent=0.25):
    """Bulleted list block."""
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, (text, fsize, fbold, fcolor) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = f"{bullet}  {text}" if bullet else text
        run.font.name  = "Calibri"
        run.font.size  = Pt(fsize)
        run.font.bold  = fbold
        run.font.color.rgb = fcolor
    return txb

def divider(slide, y, color=GOLD):
    """Thin horizontal rule."""
    line = slide.shapes.add_shape(1, Inches(0.4), Inches(y),
                                   Inches(12.53), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()

def two_col_header(slide, y, left_label, right_label,
                   lx=0.4, lw=5.9, rx=6.6, rw=6.3):
    """Two-column section headers."""
    rect(slide, lx, y, lw, 0.35, NAVY)
    rect(slide, rx, y, rw, 0.35, GREEN)
    # size bumped 11→12
    textbox(slide, lx + 0.1, y + 0.04, lw - 0.2, 0.28,
            left_label, size=12, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER)
    textbox(slide, rx + 0.1, y + 0.04, rw - 0.2, 0.28,
            right_label, size=12, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1: COVER
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()

# Full navy background
rect(slide, 0, 0, 13.33, 7.5, NAVY)

# Gold accent bars
rect(slide, 0, 2.8,  13.33, 0.06, GOLD)
rect(slide, 0, 5.55, 13.33, 0.06, GOLD)

# Main title
textbox(slide, 0.8, 1.0, 11.5, 0.8,
        "UHNW PORTFOLIO TRANSITION ANALYSIS",
        size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

textbox(slide, 0.8, 1.85, 11.5, 0.7,
        "From Concentrated Risk to Institutional Diversification",
        size=22, bold=False, color=GOLD, align=PP_ALIGN.CENTER)

# Client profile block — size bumped 11→12
textbox(slide, 2.5, 3.1, 8.0, 0.45,
        "CLIENT PROFILE",
        size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

profile_lines = [
    ("Segment:   Ultra High Net Worth (UHNW)  |  $10M+ Investable Assets",
     14, False, LGRAY),
    ("Profile:     Technology Founder  |  $25,000,000 Portfolio  |  90% Concentrated Equity",
     14, False, LGRAY),
    ("Objective:  De-risk position, generate income, fund $10M Private Foundation",
     14, False, LGRAY),
]
multiline_textbox(slide, 1.5, 3.55, 10.3, 1.6, profile_lines,
                  align=PP_ALIGN.LEFT)

# Stats row
stats = [
    ("$25M",       "Portfolio Value"),
    ("$18.5M",     "Unrealized Gains"),
    ("90%",        "Tech Concentration"),
    ("$10M",       "Foundation Goal"),
]
sw = 2.8
for i, (val, lbl) in enumerate(stats):
    sx = 0.6 + i * (sw + 0.15)
    rect(slide, sx, 4.9, sw, 0.98, RGBColor(0x00, 0x1A, 0x4A))
    textbox(slide, sx + 0.1, 4.95, sw - 0.2, 0.5,
            val, size=26, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    # label bumped 10→11
    textbox(slide, sx + 0.1, 5.47, sw - 0.2, 0.35,
            lbl, size=11, color=MGRAY, align=PP_ALIGN.CENTER)

# Footer
textbox(slide, 0.4, 6.65, 12.5, 0.35,
        "Case Study  |  February 2026  |  For Illustrative Purposes Only",
        size=10, color=MGRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2: AGENDA
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
rect(slide, 0, 0, 13.33, 7.5, LGRAY)
slide_header(slide, "AGENDA", "What we will cover today")
nav_bar(slide, 2)

agenda_items = [
    ("01", "The Challenge",
     "Why a 90% concentrated tech position represents a structural risk — not just volatility."),
    ("02", "Risk by the Numbers",
     "VaR, maximum drawdown, beta, and correlation data quantifying the exposure."),
    ("03", "The Recommendation",
     "A 7-asset-class institutional model targeting better risk-adjusted returns."),
    ("04", "Asset Allocation",
     "Side-by-side breakdown: current portfolio vs. optimized diversified model."),
    ("05", "Tax Transition Strategy",
     "Managing the $6.86M tax liability through phased liquidation and offset strategies."),
    ("06", "Securities-Based Lending",
     "Using a $5M credit facility to seed the $10M Private Foundation — without triggering additional gains."),
    ("07", "Goals-Based Planning",
     "20-year projection across bear, base, and bull scenarios."),
    ("08", "Implementation Timeline",
     "18-month phased roadmap with specialist team and key milestones."),
]

col_break = 4
# item_h bumped 0.66→0.80 to give desc text room at larger font size
item_h    = 0.80
for i, (num, title, desc) in enumerate(agenda_items):
    col = i // col_break
    row = i %  col_break
    x   = 0.4  + col * 6.5
    y   = 1.35 + row * item_h

    rect(slide, x, y, 6.1, item_h - 0.08, WHITE)
    # Number accent
    rect(slide, x, y, 0.38, item_h - 0.08, NAVY)
    # size bumped 12→13
    textbox(slide, x + 0.02, y + 0.14, 0.34, 0.40,
            num, size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    # title bumped 12→13
    textbox(slide, x + 0.45, y + 0.06, 5.5, 0.30,
            title, size=13, bold=True, color=NAVY)
    # desc bumped 9→10; y offset adjusted to avoid title overlap
    textbox(slide, x + 0.45, y + 0.39, 5.5, 0.32,
            desc,  size=10, color=DGRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3: THE CHALLENGE
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
rect(slide, 0, 0, 13.33, 7.5, LGRAY)
slide_header(slide, "THE CHALLENGE",
             "Concentrated equity is not a portfolio — it is a bet.")
nav_bar(slide, 3)

# Left: problem statement
rect(slide, 0.4, 1.35, 5.9, 5.55, WHITE)
# bumped 10→11
textbox(slide, 0.6, 1.45, 5.5, 0.35,
        "CURRENT STATE", size=11, bold=True, color=NAVY)
divider(slide, 1.84)

problems = [
    ("90% Allocated to Tech Equity",
     "A single sector drives virtually all portfolio outcomes. One regulatory shift, "
     "earnings miss, or macro rotation can permanently impair wealth."),
    ("$18.5M in Unrealized Gains",
     "The position cannot be unwound without a significant tax event. "
     "Inaction is not neutral — it compounds the concentration problem."),
    ("45% Maximum Drawdown Exposure",
     "In a 2008- or 2022-style drawdown, the portfolio could lose nearly half its value "
     "with no fixed income or alternatives to cushion the decline."),
    ("$112,500 Annual Income",
     "A 0.5% dividend yield on a $25M portfolio generates minimal cash flow. "
     "Lifestyle and philanthropic goals require a structural income solution."),
    ("No Alternative Asset Exposure",
     "Zero allocation to private equity, real estate, or hedge strategies. "
     "The client captures none of the illiquidity premium available at this wealth level."),
]
y_offset = 1.97
for title, desc in problems:
    rect(slide, 0.5, y_offset, 0.08, 0.55, RED)
    # title bumped 10→11
    textbox(slide, 0.65, y_offset, 5.55, 0.30,
            title, size=11, bold=True, color=RED)
    # desc bumped 9→10
    textbox(slide, 0.65, y_offset + 0.30, 5.55, 0.48,
            desc, size=10, color=DGRAY)
    y_offset += 1.00

# Right: visual risk summary
rect(slide, 6.6, 1.35, 6.3, 5.55, NAVY)
# bumped 10→11
textbox(slide, 6.8, 1.45, 5.9, 0.35,
        "THE COST OF INACTION", size=11, bold=True, color=GOLD)

risk_stats = [
    ("Portfolio Beta",        "1.35",    "vs. 0.88 optimized",   RED),
    ("Correlation to S&P",   "0.94",    "Almost no true diversification", RED),
    ("Value at Risk (95%)",  "$4.25M",  "Single-day loss potential",       RED),
    ("Max Drawdown (est.)",  "45%",     "Historical analog: 2022 tech rout", RED),
    ("Sharpe Ratio",         "0.38",    "Poor risk-adjusted return",        RED),
    ("Income Generated",     "$112K/yr","vs. $625K in optimized model",    RED),
]
y_r = 1.97
for label, val, sub, col in risk_stats:
    rect(slide, 6.7, y_r, 6.1, 0.78, RGBColor(0x00, 0x1A, 0x4A))
    # label bumped 9→10
    textbox(slide, 6.85, y_r + 0.05, 3.0, 0.30,
            label, size=10, color=MGRAY)
    # val bumped 16→17
    textbox(slide, 9.9, y_r + 0.03, 2.8, 0.38,
            val, size=17, bold=True, color=col, align=PP_ALIGN.RIGHT)
    # sub bumped 8→9
    textbox(slide, 6.85, y_r + 0.48, 5.8, 0.26,
            sub, size=9, color=MGRAY)
    y_r += 0.86

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4: RISK BY THE NUMBERS
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
rect(slide, 0, 0, 13.33, 7.5, LGRAY)
slide_header(slide, "RISK BY THE NUMBERS",
             "Quantifying the exposure before recommending a solution.")
nav_bar(slide, 4)

# 6 KPI boxes: top row current, bottom row optimized
kpi_data = [
    # label,               current,    optimized,   improvement, good_direction_down
    ("Portfolio Beta",       "1.35",    "0.88",    "↓ 34.8%",  True),
    ("Value at Risk (95%)", "$4.25M",  "$1.875M", "↓ 55.9%",  True),
    ("Max Drawdown (est.)", "45.0%",   "22.0%",   "↓ 51.1%",  True),
    ("Concentration Risk",  "95/100",  "15/100",  "↓ 84.2%",  True),
    ("Sharpe Ratio",        "0.38",    "0.68",    "↑ 78.9%",  False),
    ("Annual Income",       "$112K",   "$625K",   "↑ 455%",   False),
]

box_w = 2.0
box_h = 2.4
gap   = 0.12
start_x = (13.33 - (6 * box_w + 5 * gap)) / 2

for i, (label, curr, opt, chg, down_is_good) in enumerate(kpi_data):
    x = start_x + i * (box_w + gap)

    # Current (top, dark red band)
    rect(slide, x, 1.35, box_w, 1.08, RGBColor(0x3A, 0x00, 0x00))
    # "CURRENT" bumped 7→9
    textbox(slide, x + 0.05, 1.38, box_w - 0.1, 0.24,
            "CURRENT", size=9, bold=True, color=RED, align=PP_ALIGN.CENTER)
    # value bumped 20→21
    textbox(slide, x + 0.05, 1.62, box_w - 0.1, 0.50,
            curr, size=21, bold=True, color=RED, align=PP_ALIGN.CENTER)
    # label bumped 8→9
    textbox(slide, x + 0.05, 2.14, box_w - 0.1, 0.24,
            label, size=9, color=MGRAY, align=PP_ALIGN.CENTER)

    # Arrow + change — bumped 11→12
    rect(slide, x, 2.45, box_w, 0.30, LGRAY)
    textbox(slide, x + 0.05, 2.47, box_w - 0.1, 0.26,
            chg, size=12, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    # Optimized (bottom, dark green band)
    rect(slide, x, 2.77, box_w, 1.08, RGBColor(0x00, 0x25, 0x15))
    # "OPTIMIZED" bumped 7→9
    textbox(slide, x + 0.05, 2.80, box_w - 0.1, 0.24,
            "OPTIMIZED", size=9, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    # value bumped 20→21
    textbox(slide, x + 0.05, 3.04, box_w - 0.1, 0.50,
            opt, size=21, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    # label bumped 8→9
    textbox(slide, x + 0.05, 3.56, box_w - 0.1, 0.24,
            label, size=9, color=MGRAY, align=PP_ALIGN.CENTER)

# Insight callout — bumped 9→10 and 10→11
rect(slide, 0.4, 4.05, 12.53, 0.92, NAVY)
rect(slide, 0.4, 4.05, 0.08,  0.92, GOLD)
textbox(slide, 0.6, 4.10, 12.2, 0.28,
        "KEY INSIGHT", size=10, bold=True, color=GOLD)
textbox(slide, 0.6, 4.40, 12.2, 0.50,
        "The current portfolio takes on 53% more market risk than the optimized model while generating "
        "82% less annual income. This is not a trade-off — it is an inefficiency. "
        "Diversification does not sacrifice return; it improves the quality of each unit of risk taken.",
        size=11, color=WHITE)

# Bottom comparison table header — bumped 9→10
rect(slide, 0.4, 5.15, 12.53, 0.32, NAVY)
for col_x, label in [(0.5, "METRIC"), (4.5, "CURRENT"), (7.5, "OPTIMIZED"), (10.5, "CHANGE")]:
    textbox(slide, col_x, 5.18, 2.8, 0.25,
            label, size=10, bold=True, color=WHITE)

table_rows = [
    ("Number of Holdings",      "8",      "145",     "+1,713%"),
    ("Asset Classes",           "2",      "7",       "+5 classes"),
    ("Geographic Exposure",     "1 Region","4 Regions","+300%"),
    ("Correlation to S&P 500",  "0.94",   "0.68",   "↓ 27.7%"),
]
for i, (metric, curr, opt, chg) in enumerate(table_rows):
    y  = 5.50 + i * 0.38
    bg = WHITE if i % 2 == 0 else LGRAY
    rect(slide, 0.4, y, 12.53, 0.36, bg)
    # all row text bumped 9→10
    textbox(slide, 0.5, y + 0.05, 3.8, 0.28, metric, size=10, color=DGRAY, bold=True)
    textbox(slide, 4.5, y + 0.05, 2.8, 0.28, curr,   size=10, color=RED,   align=PP_ALIGN.CENTER)
    textbox(slide, 7.5, y + 0.05, 2.8, 0.28, opt,    size=10, color=GREEN, align=PP_ALIGN.CENTER)
    textbox(slide, 10.5,y + 0.05, 2.5, 0.28, chg,    size=10, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5: THE RECOMMENDATION
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
rect(slide, 0, 0, 13.33, 7.5, LGRAY)
slide_header(slide, "THE RECOMMENDATION",
             "An institutional-grade model built for long-term wealth preservation and growth.")
nav_bar(slide, 5)

# 3-column framework
cols = [
    ("DIVERSIFY",    NAVY,  "Expand from 2 to 7 asset classes across public equity, "
                            "fixed income, private equity, real estate, and cash."),
    ("DE-RISK",      GREEN, "Reduce portfolio beta from 1.35 to 0.88. Cut Value at Risk "
                            "by 55.9%. Target a Sharpe Ratio improvement of 78.9%."),
    ("GENERATE",     GOLD,  "Grow annual income from $112,500 to $625,000. "
                            "Use $5M SBL credit facility to seed $10M Private Foundation "
                            "without triggering additional capital gains."),
]
cw2 = 4.0
for i, (title, color, desc) in enumerate(cols):
    cx = 0.4 + i * (cw2 + 0.27)
    rect(slide, cx, 1.35, cw2, 0.52, color)
    # title bumped 18→19
    textbox(slide, cx + 0.1, 1.40, cw2 - 0.2, 0.42,
            title, size=19, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(slide, cx, 1.87, cw2, 3.65, WHITE)
    # desc bumped 11→12; h kept at 1.0 — clears TARGET ALLOCATION label (now at y=2.77+)
    textbox(slide, cx + 0.15, 1.97, cw2 - 0.3, 0.90,
            desc, size=12, color=DGRAY)

# Model allocations
# FIX: y_a moved from 3.05 to 3.35 to eliminate overlap with column desc textboxes
# (desc boxes end at y≈2.87, label now starts at y=3.07 — 0.20" clear)
alloc_data = [
    ("Public Equity",      "55%", "US, International, EM across large, mid, small cap"),
    ("Fixed Income",       "23%", "Investment grade, high yield, municipal bonds"),
    ("Private Equity",     "12%", "Buyout funds + venture capital — illiquidity premium"),
    ("Real Estate",        "8%",  "Core + opportunistic — inflation hedge, income"),
    ("Cash & Equivalents", "2%",  "Money market — liquidity buffer"),
]
y_a = 3.35
# bumped 9→10
textbox(slide, 0.55, y_a - 0.26, 12.0, 0.26,
        "TARGET ALLOCATION  |  OPTIMIZED DIVERSIFIED MODEL",
        size=10, bold=True, color=NAVY)
rect(slide, 0.4, y_a - 0.03, 12.53, 0.32, NAVY)
# bumped 8→9
for col_x, lbl in [(0.5,"ASSET CLASS"),(4.8,"ALLOCATION"),(6.9,"DESCRIPTION")]:
    textbox(slide, col_x, y_a - 0.00, 2.2, 0.28,
            lbl, size=9, bold=True, color=WHITE)

bar_colors = [NAVY, GREEN, RGBColor(0x70, 0x30, 0xA0),
              RGBColor(0xC8, 0x50, 0x00), MGRAY]

for i, (ac, alloc, desc) in enumerate(alloc_data):
    y_r = y_a + 0.32 + i * 0.52
    bg  = WHITE if i % 2 == 0 else LGRAY
    rect(slide, 0.4,  y_r, 12.53, 0.50, bg)
    rect(slide, 0.4,  y_r, 0.12,  0.50, bar_colors[i])
    # ac name bumped 10→11
    textbox(slide, 0.6, y_r + 0.09, 4.0, 0.32,
            ac, size=11, bold=True, color=DGRAY)
    # alloc pct bumped 16→17
    textbox(slide, 4.8, y_r + 0.07, 1.8, 0.36,
            alloc, size=17, bold=True, color=bar_colors[i],
            align=PP_ALIGN.CENTER)
    # desc bumped 9→10
    textbox(slide, 6.9, y_r + 0.09, 5.8, 0.32,
            desc, size=10, color=DGRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6: ASSET ALLOCATION
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
rect(slide, 0, 0, 13.33, 7.5, LGRAY)
slide_header(slide, "ASSET ALLOCATION",
             "Before and after: a side-by-side comparison of the full balance sheet.")
nav_bar(slide, 6)

two_col_header(slide, 1.35, "CURRENT PORTFOLIO", "OPTIMIZED DIVERSIFIED MODEL")

# Current portfolio breakdown
current_allocs = [
    ("US Large Cap Tech",  90.0, RED),
    ("Cash & Equivalents", 10.0, MGRAY),
]
# Optimized portfolio breakdown (simplified display)
optimized_display = [
    ("US Large Cap",          25.0, NAVY),
    ("International / EM",    17.0, RGBColor(0x1F, 0x5C, 0x99)),
    ("Fixed Income",          23.0, GREEN),
    ("Private Equity",        12.0, RGBColor(0x70, 0x30, 0xA0)),
    ("Real Estate",            8.0, RGBColor(0xC8, 0x50, 0x00)),
    ("US Mid & Small Cap",    13.0, RGBColor(0x00, 0x4A, 0x80)),
    ("Cash & Equivalents",     2.0, MGRAY),
]

def bar_chart(slide, x, y, w, h, allocs, total=100.0):
    """Horizontal stacked bar chart."""
    bar_h = 0.65
    drawn = 0.0
    for label, pct, color in allocs:
        bw = (pct / total) * w
        rect(slide, x + (drawn / total) * w, y, bw, bar_h, color)
        drawn += pct
    # Labels below — bumped 7→9
    drawn = 0.0
    for label, pct, color in allocs:
        bw = (pct / total) * w
        lx = x + (drawn / total) * w
        if bw > 0.6:
            textbox(slide, lx + 0.03, y + bar_h + 0.05, bw - 0.06, 0.36,
                    f"{label}\n{pct:.0f}%", size=9, color=DGRAY)
        drawn += pct

bar_chart(slide, 0.5, 1.85, 5.7, 0.65, current_allocs)
bar_chart(slide, 6.7, 1.85, 5.9, 0.65, optimized_display)

# Current detail table — bumped 8→9 headers, 9→10 rows
y_t = 2.88
rect(slide, 0.4, y_t, 5.9, 0.32, NAVY)
for col_x, lbl in [(0.5,"HOLDING"),(3.2,"ALLOC"),(4.4,"VALUE")]:
    textbox(slide, col_x, y_t + 0.04, 1.5, 0.25,
            lbl, size=9, bold=True, color=WHITE)

curr_rows = [
    ("US Large Cap Tech",    "90.0%", "$22,500,000"),
    ("Cash & Equivalents",   "10.0%", " $2,500,000"),
    ("",                     "",      ""),
    ("TOTAL",                "100%",  "$25,000,000"),
]
for i, (nm, al, vl) in enumerate(curr_rows):
    yr = y_t + 0.32 + i * 0.44
    bg = WHITE if i % 2 == 0 else LGRAY
    bold = nm == "TOTAL"
    bg2  = NAVY if bold else bg
    rect(slide, 0.4, yr, 5.9, 0.42, bg2)
    fc   = WHITE if bold else DGRAY
    textbox(slide, 0.5,  yr+0.08, 2.5,  0.28, nm, size=10, bold=bold, color=fc)
    textbox(slide, 3.2,  yr+0.08, 0.9,  0.28, al, size=10, bold=bold, color=fc,
            align=PP_ALIGN.CENTER)
    textbox(slide, 4.2,  yr+0.08, 1.9,  0.28, vl, size=10, bold=bold, color=fc,
            align=PP_ALIGN.RIGHT)

# Optimized detail table — bumped 8→9 headers, 8→9 rows
y_t2 = 2.88
rect(slide, 6.6, y_t2, 6.3, 0.32, GREEN)
for col_x, lbl in [(6.7,"ASSET CLASS"),(9.8,"ALLOC"),(11.1,"VALUE")]:
    textbox(slide, col_x, y_t2 + 0.04, 1.8, 0.25,
            lbl, size=9, bold=True, color=WHITE)

opt_rows = [
    ("Public Equity (Total)",      "55%",  "$13,750,000"),
    ("  — US Large Cap",           "25%",  " $6,250,000"),
    ("  — Intl / EM",              "17%",  " $4,250,000"),
    ("  — US Mid & Small Cap",     "13%",  " $3,250,000"),
    ("Fixed Income",               "23%",  " $5,750,000"),
    ("Private Equity",             "12%",  " $3,000,000"),
    ("Real Estate",                 "8%",  " $2,000,000"),
    ("Cash & Equivalents",          "2%",  "   $500,000"),
    ("TOTAL",                      "100%", "$25,000,000"),
]
for i, (nm, al, vl) in enumerate(opt_rows):
    yr   = y_t2 + 0.32 + i * 0.38
    bold = nm == "TOTAL"
    sub  = nm.startswith("  —")
    bg2  = NAVY if bold else (LGRAY if i % 2 == 0 else WHITE)
    fc   = WHITE if bold else (MGRAY if sub else DGRAY)
    rect(slide, 6.6, yr, 6.3, 0.36, bg2)
    textbox(slide, 6.7,  yr+0.06, 2.9, 0.26, nm, size=9, bold=bold, color=fc)
    textbox(slide, 9.8,  yr+0.06, 0.9, 0.26, al, size=9, bold=bold, color=fc,
            align=PP_ALIGN.CENTER)
    textbox(slide, 11.0, yr+0.06, 1.7, 0.26, vl, size=9, bold=bold, color=fc,
            align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7: TAX TRANSITION STRATEGY
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
rect(slide, 0, 0, 13.33, 7.5, LGRAY)
slide_header(slide, "TAX TRANSITION STRATEGY",
             "The $6.86M question — and five ways to reduce it.")
nav_bar(slide, 7)

# Left: the tax problem
rect(slide, 0.4, 1.35, 5.9, 5.55, WHITE)
# bumped 10→11
textbox(slide, 0.6, 1.45, 5.5, 0.30,
        "THE TAX REALITY", size=11, bold=True, color=NAVY)
divider(slide, 1.80)

tax_facts = [
    ("$18,500,000",   "Unrealized Capital Gains",         RED),
    ("37.1%",         "Blended Effective Rate (Fed + NIIT + CA)", RED),
    ("$6,863,500",    "Estimated Gross Tax Liability",     RED),
    ("($850,000)",    "Tax-Loss Harvesting Offset",        GREEN),
    ("$6,013,500",    "NET Tax Liability After Offset",    RED),
]
y_tf = 1.97
for val, label, color in tax_facts:
    is_net = label == "NET Tax Liability After Offset"
    if is_net:
        rect(slide, 0.5, y_tf, 5.7, 0.62, LGRAY)
    # tax values bumped 18→19
    textbox(slide, 0.6, y_tf + 0.02, 2.2, 0.40,
            val, size=19, bold=True, color=color, align=PP_ALIGN.RIGHT)
    # labels bumped 9→10
    textbox(slide, 2.9, y_tf + 0.11, 3.2, 0.30,
            label, size=10, color=DGRAY, bold=is_net)
    y_tf += 0.70

# Phased liquidation mini-chart — label bumped 8→9
textbox(slide, 0.6, y_tf + 0.05, 5.5, 0.28,
        "5-TRANCHE LIQUIDATION SCHEDULE", size=9, bold=True, color=NAVY)
tranches = [("T1\nMo.4","$4.5M"),("T2\nMo.5","$4.5M"),
            ("T3\nMo.6","$4.5M"),("T4\nMo.7","$4.5M"),("T5\nMo.8","$4.5M")]
tw = 0.95
for j, (lbl, amt) in enumerate(tranches):
    tx = 0.5 + j * (tw + 0.12)
    rect(slide, tx, y_tf + 0.38, tw, 0.65, NAVY)
    # tranche labels bumped 7→8
    textbox(slide, tx, y_tf + 0.40, tw, 0.30,
            lbl, size=8, color=GOLD, align=PP_ALIGN.CENTER)
    # amounts bumped 9→10
    textbox(slide, tx, y_tf + 0.69, tw, 0.28,
            amt, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Right: 5 offset strategies
rect(slide, 6.6, 1.35, 6.3, 5.55, NAVY)
# bumped 10→11
textbox(slide, 6.8, 1.45, 5.9, 0.30,
        "TAX MITIGATION STRATEGIES", size=11, bold=True, color=GOLD)

strategies = [
    ("Tax-Loss Harvesting",
     "$850K offset from identified losses across portfolio positions."),
    ("Installment Liquidation",
     "5 monthly tranches (Months 4-8). Straddle December/January to split gains across two tax years."),
    ("Donor-Advised Fund (DAF)",
     "Donate appreciated shares directly. Avoid gains entirely; receive FMV deduction."),
    ("Qualified Opportunity Zone",
     "Invest deferred gains in QOZ fund. Defer recognition; exclude future appreciation."),
    ("Charitable Remainder Trust",
     "Transfer stock to CRT. Avoid immediate gain; receive income stream and estate benefit."),
]
y_s = 1.90
for i, (title, desc) in enumerate(strategies):
    rect(slide, 6.7, y_s, 6.1, 0.90, RGBColor(0x00, 0x1A, 0x4A))
    rect(slide, 6.7, y_s, 0.06, 0.90, GOLD)
    # strategy titles bumped 10→11
    textbox(slide, 6.85, y_s + 0.06, 5.8, 0.30,
            f"{i+1}.  {title}", size=11, bold=True, color=WHITE)
    # strategy descs bumped 9→10
    textbox(slide, 6.85, y_s + 0.38, 5.8, 0.46,
            desc, size=10, color=MGRAY)
    y_s += 1.02

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8: SECURITIES-BASED LENDING
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
rect(slide, 0, 0, 13.33, 7.5, LGRAY)
slide_header(slide, "SECURITIES-BASED LENDING (SBL)",
             "Seed the $10M Private Foundation without selling a single share.")
nav_bar(slide, 8)

# Top concept bar — bumped 9→10 and 10→11
rect(slide, 0.4, 1.35, 12.53, 0.72, NAVY)
rect(slide, 0.4, 1.35, 0.08,  0.72, GOLD)
textbox(slide, 0.6, 1.38, 12.2, 0.28,
        "THE CORE INSIGHT", size=10, bold=True, color=GOLD)
textbox(slide, 0.6, 1.65, 12.2, 0.38,
        "The tax cost of liquidating $5M of tech stock to fund the foundation: $1,669,500. "
        "The annual after-tax cost of a $5M SBL credit facility at 6.50%: $204,750. "
        "Over 10 years, the credit strategy outperforms liquidation by $7,124,500 — "
        "keeping $5M compounding in the portfolio while the foundation receives full seed capital.",
        size=11, color=WHITE)

# Loan structure — header bumped 9→10
rect(slide, 0.4, 2.22, 5.9, 0.32, GREEN)
textbox(slide, 0.5, 2.26, 5.7, 0.26,
        "CREDIT FACILITY STRUCTURE", size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

loan_items = [
    ("Loan Amount",              "$5,000,000"),
    ("Collateral Pledged",       "$15,000,000  (60% of portfolio)"),
    ("Loan-to-Value (LTV)",      "33.3%  —  conservative"),
    ("Margin Call Trigger",      "40.0% LTV"),
    ("Interest Rate",            "6.50%  (SOFR + 200 bps)"),
    ("After-Tax Effective Rate", "4.10%  (at 37% bracket)"),
    ("Annual Interest Cost",     "$325,000 gross  /  $204,750 net"),
    ("Loan Structure",           "Interest-only  |  Annually renewable"),
    ("Use of Proceeds",          "Private Foundation seed capital"),
]
y_l = 2.58
for i, (param, val) in enumerate(loan_items):
    bg = WHITE if i % 2 == 0 else LGRAY
    bold = param in ["Loan Amount", "After-Tax Effective Rate", "Annual Interest Cost"]
    rect(slide, 0.4, y_l, 5.9, 0.42, bg)
    # bumped 9→10
    textbox(slide, 0.5, y_l + 0.08, 2.6, 0.28,
            param, size=10, bold=bold, color=NAVY if bold else DGRAY)
    textbox(slide, 3.2, y_l + 0.08, 2.9, 0.28,
            val,   size=10, bold=bold, color=GREEN if bold else DGRAY,
            align=PP_ALIGN.RIGHT)
    y_l += 0.42

# Right: SBL vs. Liquidation comparison — bumped 9→10
rect(slide, 6.6, 2.22, 6.3, 0.32, NAVY)
textbox(slide, 6.7, 2.26, 6.1, 0.26,
        "SBL vs. LIQUIDATION — 10-YEAR COMPARISON",
        size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

compare_rows = [
    ("Factor",                  "Liquidate",    "Use SBL",      True),
    ("Immediate Tax Cost",      "$1,669,500",   "$0",           False),
    ("Annual Carry Cost",       "$0",           "$204,750",     False),
    ("Capital to Foundation",   "$3,330,500",   "$5,000,000",   False),
    ("Portfolio Compounding",   "$20M base",    "$25M base",    False),
    ("10-Yr Interest Cost",     "$0",           "$2,047,500",   False),
    ("Net 10-Yr SBL Advantage", "",             "$7,124,500",   False),
]
y_c = 2.58
for i, (factor, liq, sbl_v, is_header) in enumerate(compare_rows):
    bg   = NAVY if is_header else (WHITE if i % 2 == 0 else LGRAY)
    fc   = WHITE if is_header else DGRAY
    rect(slide, 6.6, y_c, 6.3, 0.42, bg)
    textbox(slide, 6.7,  y_c + 0.08, 2.4, 0.28, factor, size=10,
            bold=is_header or factor == "Net 10-Yr SBL Advantage",
            color=fc if not is_header else WHITE)
    textbox(slide, 9.2,  y_c + 0.08, 1.7, 0.28, liq, size=10,
            color=WHITE if is_header else RED, align=PP_ALIGN.CENTER)
    textbox(slide, 11.1, y_c + 0.08, 1.7, 0.28, sbl_v, size=10, bold=not is_header,
            color=WHITE if is_header else GREEN, align=PP_ALIGN.CENTER)
    y_c += 0.42

# Margin call stress test
# FIX: Moved header up to y=5.55 (was 5.60) and reduced row height to 0.28 (was 0.30)
# so the last row bottom = 5.83 + 3*0.28 + 0.28 = 6.95 < 7.1 (nav bar). No more overlap.
rect(slide, 6.6, 5.55, 6.3, 0.28, NAVY)
# bumped 8→9
textbox(slide, 6.7, 5.57, 6.1, 0.22,
        "MARGIN CALL STRESS TEST  (Collateral: $15M  |  Loan: $5M)",
        size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
stress_rows = [
    ("Base Case  (0% decline)",  "33.3%", "NO"),
    ("Moderate   (-10%)",        "37.0%", "NO"),
    ("Significant (-20%)",       "41.7%", "YES ⚠"),
    ("Severe     (-30%)",        "47.6%", "YES ⚠"),
]
y_st = 5.83   # was 5.90 — shifted up to keep content above nav bar
for i, (scenario, ltv, call) in enumerate(stress_rows):
    bg  = WHITE if i % 2 == 0 else LGRAY
    col = RED if "YES" in call else GREEN
    rect(slide, 6.6, y_st, 6.3, 0.28, bg)
    # bumped 8→9
    textbox(slide, 6.7,  y_st + 0.04, 3.2, 0.22, scenario, size=9, color=DGRAY)
    textbox(slide, 9.9,  y_st + 0.04, 1.2, 0.22, f"LTV: {ltv}", size=9,
            color=RED if "YES" in call else DGRAY, align=PP_ALIGN.CENTER)
    textbox(slide, 11.3, y_st + 0.04, 1.4, 0.22, call, size=9, bold=True,
            color=col, align=PP_ALIGN.CENTER)
    y_st += 0.28   # was 0.30

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9: GOALS-BASED PLANNING
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
rect(slide, 0, 0, 13.33, 7.5, LGRAY)
slide_header(slide, "GOALS-BASED PLANNING",
             "Three scenarios. One consistent conclusion: the optimized model funds your goals.")
nav_bar(slide, 9)

# Scenario assumption boxes
scenarios_info = [
    ("BEAR CASE",  "5.0% Return",  "Net 3.0% after 2% withdrawal", RED,   "Conservative markets, elevated volatility"),
    ("BASE CASE",  "9.8% Return",  "Net 7.8% after 2% withdrawal", NAVY,  "Institutional long-term capital market assumptions"),
    ("BULL CASE",  "13.5% Return", "Net 11.5% after 2% withdrawal",GREEN, "Favorable market environment, alternatives outperform"),
]
bw2 = 3.9
for i, (title, ret, net, color, note) in enumerate(scenarios_info):
    sx = 0.4 + i * (bw2 + 0.22)
    rect(slide, sx, 1.35, bw2, 0.42, color)
    # title bumped 11→12
    textbox(slide, sx + 0.1, 1.38, bw2 - 0.2, 0.36,
            f"{title}  |  {ret}", size=12, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER)
    rect(slide, sx, 1.77, bw2, 0.50, WHITE)
    # net bumped 9→10
    textbox(slide, sx + 0.1, 1.80, bw2 - 0.2, 0.22,
            net, size=10, bold=True, color=color, align=PP_ALIGN.CENTER)
    # note bumped 8→9
    textbox(slide, sx + 0.1, 2.03, bw2 - 0.2, 0.22,
            note, size=9, color=DGRAY, align=PP_ALIGN.CENTER)

# 20-year projection table — headers bumped 8→9, rows bumped 9→10
rect(slide, 0.4, 2.44, 12.53, 0.30, NAVY)
for col_x, lbl in [(0.5,"YEAR"),(2.0,"CALENDAR"),(4.0,"BEAR (3.0% net)"),
                    (7.0,"BASE (7.8% net)"),(10.0,"BULL (11.5% net)")]:
    textbox(slide, col_x, 2.47, 2.5, 0.24,
            lbl, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

start = 18_986_500
bear_v = base_v = bull_v = start
bear_r, base_r, bull_r = 0.030, 0.078, 0.115

display_years = [1, 3, 5, 7, 10, 15, 20]
milestones    = {
    1:  "SBL Established",
    5:  "5-Year Review",
    7:  "Foundation Review",
    10: "Loan Repaid",
    20: "Planning Horizon",
}

# Run projection
all_vals = {}
bv = base_v
bev = bear_v
buv = bull_v
for yr in range(1, 21):
    bev *= (1 + bear_r)
    bv  *= (1 + base_r)
    buv *= (1 + bull_r)
    all_vals[yr] = (bev, bv, buv)

y_p = 2.76
for i, yr in enumerate(display_years):
    bev, bv, buv = all_vals[yr]
    bg   = WHITE if i % 2 == 0 else LGRAY
    bold = yr in [10, 20]
    rect(slide, 0.4, y_p, 12.53, 0.42, bg)
    # all row values bumped 9→10
    textbox(slide, 0.5,  y_p + 0.09, 1.3,  0.28,
            f"Year {yr:>2}", size=10, bold=bold, color=DGRAY)
    textbox(slide, 2.0,  y_p + 0.09, 1.6,  0.28,
            str(2026 + yr), size=10, bold=bold, color=DGRAY,
            align=PP_ALIGN.CENTER)
    textbox(slide, 4.0,  y_p + 0.09, 2.5,  0.28,
            f"${bev:,.0f}", size=10, bold=bold, color=RED if bold else DGRAY,
            align=PP_ALIGN.CENTER)
    textbox(slide, 7.0,  y_p + 0.09, 2.5,  0.28,
            f"${bv:,.0f}",  size=10, bold=bold, color=NAVY if bold else DGRAY,
            align=PP_ALIGN.CENTER)
    textbox(slide, 10.0, y_p + 0.09, 2.2,  0.28,
            f"${buv:,.0f}", size=10, bold=bold, color=GREEN if bold else DGRAY,
            align=PP_ALIGN.CENTER)
    if yr in milestones:
        # milestone bumped 7→9
        textbox(slide, 12.3, y_p + 0.09, 0.98, 0.28,
                milestones[yr], size=9, color=RGBColor(0x70, 0x30, 0xA0))
    y_p += 0.42

# Foundation callout — bumped 9→10
rect(slide, 0.4, 5.75, 12.53, 0.85, NAVY)
rect(slide, 0.4, 5.75, 0.08,  0.85, GOLD)
textbox(slide, 0.6, 5.79, 12.2, 0.28,
        "PRIVATE FOUNDATION MILESTONE", size=10, bold=True, color=GOLD)
textbox(slide, 0.6, 6.08, 12.2, 0.46,
        "$5M deployed via SBL in Year 1 (no liquidation). Foundation compounds at 2.5% net (7.5% return less 5% annual grantmaking). "
        "$70,760/yr directed from diversified portfolio income ($625K/yr generated) bridges seed to $10M target. "
        "Projected value: $7.2M at Year 10, $10.0M at Year 20. Loan repaid from portfolio income — no asset sales required.",
        size=10, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10: IMPLEMENTATION TIMELINE
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
rect(slide, 0, 0, 13.33, 7.5, LGRAY)
slide_header(slide, "IMPLEMENTATION TIMELINE",
             "An 18-month phased roadmap with a specialist team at every stage.")
nav_bar(slide, 10)

phases = [
    ("DISCOVERY\nMonths 1-3",   NAVY,  [
        "Complete Risk Tolerance Questionnaire",
        "Review Holdings & Cost Basis",
        "Tax Consultation with CPA",
        "Develop Investment Policy Statement",
        "Model Transition Scenarios",
    ]),
    ("EXECUTION\nMonths 4-11",  GREEN, [
        "5-Tranche Systematic Liquidation",
        "Deploy to Fixed Income ($5.75M)",
        "Public Equity Diversification",
        "Private Equity Capital Calls",
        "Real Estate Investment Deployment",
    ]),
    ("MONITORING\nMonths 12-18", RGBColor(0x70, 0x30, 0xA0), [
        "Portfolio Rebalancing Review",
        "Tax-Loss Harvesting ($850K target)",
        "SBL Facility Establishment",
        "Quarterly Performance Reviews",
        "Semi-Annual Risk Assessment",
    ]),
]

pw = 4.0
for i, (title, color, items) in enumerate(phases):
    px = 0.4 + i * (pw + 0.27)
    rect(slide, px, 1.35, pw, 0.62, color)
    # title bumped 12→13
    textbox(slide, px + 0.1, 1.37, pw - 0.2, 0.58,
            title, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(slide, px, 1.97, pw, 4.85, WHITE)
    for j, item in enumerate(items):
        iy = 2.05 + j * 0.88
        rect(slide, px + 0.1, iy, 0.06, 0.60, color)
        # item text bumped 10→11
        textbox(slide, px + 0.22, iy + 0.08, pw - 0.4, 0.30,
                item, size=11, bold=True, color=DGRAY)
        party_map = {
            "Complete Risk Tolerance Questionnaire": "Private Banker",
            "Review Holdings & Cost Basis":          "Investment Analyst",
            "Tax Consultation with CPA":             "Wealth Planning Specialist",
            "Develop Investment Policy Statement":   "Private Banker",
            "Model Transition Scenarios":            "Investment Analyst",
            "5-Tranche Systematic Liquidation":      "Trading Desk",
            "Deploy to Fixed Income ($5.75M)":       "Fixed Income Specialist",
            "Public Equity Diversification":         "Equity Specialist",
            "Private Equity Capital Calls":          "Alternatives Specialist",
            "Real Estate Investment Deployment":     "Real Estate Specialist",
            "Portfolio Rebalancing Review":          "Private Banker",
            "Tax-Loss Harvesting ($850K target)":   "Wealth Planning Specialist",
            "SBL Facility Establishment":            "Credit / Banking",
            "Quarterly Performance Reviews":         "Private Banker",
            "Semi-Annual Risk Assessment":           "Investment Analyst",
        }
        party = party_map.get(item, "")
        # party bumped 8→9
        textbox(slide, px + 0.22, iy + 0.42, pw - 0.4, 0.26,
                party, size=9, color=MGRAY, italic=True)

# Timeline bar at bottom — bumped 7→9
rect(slide, 0.4, 6.65, 12.53, 0.28, NAVY)
for month_x, label in [
    (0.4,  "Month 1"), (1.85, "Month 3"), (3.3,  "Month 4"),
    (5.85, "Month 8"), (7.3,  "Month 12"),(10.7, "Month 18"),
]:
    textbox(slide, month_x + 0.05, 6.67, 1.3, 0.22,
            label, size=9, color=GOLD, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11: SUMMARY
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
rect(slide, 0, 0, 13.33, 7.5, NAVY)
rect(slide, 0, 1.05, 13.33, 0.06, GOLD)

# bumped 28→30
textbox(slide, 0.6, 0.10, 12.0, 0.80,
        "SUMMARY & NEXT STEPS",
        size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 5 outcome boxes
outcomes_final = [
    ("Risk Reduction",    "55.9%",  "decrease in Value at Risk",          RED),
    ("Return Quality",    "+78.9%", "improvement in risk-adjusted return", GREEN),
    ("Income Growth",     "+455%",  "annual portfolio income increase",    GREEN),
    ("Tax Efficiency",    "$850K",  "tax-loss harvesting offset identified",GOLD),
    ("Foundation",        "$5M",    "deployed via SBL — zero liquidation", GREEN),
]
bw3 = 2.4
for i, (title, val, desc, color) in enumerate(outcomes_final):
    ox = 0.4 + i * (bw3 + 0.15)
    rect(slide, ox, 1.28, bw3, 1.45, RGBColor(0x00, 0x1A, 0x4A))
    # title bumped 9→10
    textbox(slide, ox + 0.1, 1.33, bw3 - 0.2, 0.30,
            title, size=10, bold=True, color=MGRAY, align=PP_ALIGN.CENTER)
    textbox(slide, ox + 0.05, 1.62, bw3 - 0.1, 0.52,
            val, size=22, bold=True, color=color, align=PP_ALIGN.CENTER)
    # desc bumped 8→9
    textbox(slide, ox + 0.05, 2.14, bw3 - 0.1, 0.52,
            desc, size=9, color=MGRAY, align=PP_ALIGN.CENTER)

# Next steps
rect(slide, 0.4, 2.95, 12.53, 0.35, GOLD)
# bumped 11→12
textbox(slide, 0.6, 2.98, 12.0, 0.28,
        "PROPOSED NEXT STEPS", size=12, bold=True, color=NAVY)

next_steps = [
    ("01", "Sign Investment Policy Statement",
     "Formalizes the agreed-upon asset allocation, risk parameters, and goals."),
    ("02", "Schedule Tax Strategy Session with CPA",
     "Coordinates the 5-tranche liquidation structure, harvesting plan, and DAF contribution."),
    ("03", "Authorize Tranche 1 Liquidation ($4.5M)",
     "Initiates the systematic de-risking process. Trading desk executes with price optimization."),
    ("04", "Establish SBL Credit Facility",
     "Secures the $5M non-purpose line against the diversified portfolio to seed the foundation."),
]
y_ns = 3.35
for i, (num, title, desc) in enumerate(next_steps):
    rect(slide, 0.4, y_ns, 12.53, 0.82, RGBColor(0x00, 0x1A, 0x4A))
    rect(slide, 0.4, y_ns, 0.45, 0.82, GOLD)
    textbox(slide, 0.45, y_ns + 0.18, 0.35, 0.40,
            num, size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    # title bumped 11→12
    textbox(slide, 0.92, y_ns + 0.06, 11.8, 0.32,
            title, size=12, bold=True, color=WHITE)
    # desc bumped 9→10
    textbox(slide, 0.92, y_ns + 0.40, 11.8, 0.36,
            desc, size=10, color=MGRAY)
    y_ns += 0.90

# Footer — bumped 7→8
rect(slide, 0, 6.98, 13.33, 0.52, RGBColor(0x00, 0x10, 0x30))
textbox(slide, 0.4, 7.02, 12.53, 0.38,
        "For illustrative purposes only. All figures are estimates based on modeled assumptions. "
        "Past performance does not guarantee future results. Consult a licensed financial advisor before making investment decisions.  "
        "|  February 2026",
        size=8, color=MGRAY, align=PP_ALIGN.CENTER)

# ── SAVE ──────────────────────────────────────────────────────────────────────
prs.save(OUTPUT_FILE)
print(f"Presentation saved: {OUTPUT_FILE}")
print(f"Slides: {len(prs.slides)}")
